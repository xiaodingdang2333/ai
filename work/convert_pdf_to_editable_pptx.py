from pathlib import Path
import io
import math

import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


PDF_PATH = Path(r"D:\ai\work\qly_dict_service.pdf")
OUT_PATH = Path(r"D:\ai\work\千里眼DICT集成服务细化管理方案宣贯材料_可编辑版.pptx")


def pdf_to_inches(v):
    return Inches(v / 72.0)


def rgb_from_pdf(color):
    if color is None:
        return None
    vals = []
    for c in color[:3]:
        vals.append(max(0, min(255, int(round(c * 255)))))
    return RGBColor(*vals)


def rgb_from_int(value):
    return RGBColor((value >> 16) & 255, (value >> 8) & 255, value & 255)


def add_rect(slide, rect, fill=None, line=None, width=0):
    x0, y0, x1, y1 = rect
    if x1 <= x0 or y1 <= y0:
        return
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        pdf_to_inches(x0),
        pdf_to_inches(y0),
        pdf_to_inches(x1 - x0),
        pdf_to_inches(y1 - y0),
    )
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is not None and width and width > 0:
        shape.line.color.rgb = line
        shape.line.width = Pt(width)
    else:
        shape.line.fill.background()


def add_line(slide, p1, p2, color, width):
    x0, y0 = p1
    x1, y1 = p2
    shape = slide.shapes.add_connector(
        1,
        pdf_to_inches(x0),
        pdf_to_inches(y0),
        pdf_to_inches(x1),
        pdf_to_inches(y1),
    )
    shape.line.color.rgb = color or RGBColor(0, 0, 0)
    shape.line.width = Pt(width or 0.5)


def add_text_line(slide, line):
    spans = [s for s in line.get("spans", []) if s.get("text")]
    if not spans:
        return
    text = "".join(s["text"] for s in spans)
    if not text.strip():
        return

    x0, y0, x1, y1 = line["bbox"]
    height = max(y1 - y0 + 2, 8)
    width = max(x1 - x0 + 4, 8)
    tx = slide.shapes.add_textbox(
        pdf_to_inches(x0),
        pdf_to_inches(y0 - 1),
        pdf_to_inches(width),
        pdf_to_inches(height),
    )
    tf = tx.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1

    for span in spans:
        run = p.add_run()
        run.text = span["text"]
        font = run.font
        font.name = "Microsoft YaHei"
        font.size = Pt(max(1, span.get("size", 10)))
        font.bold = bool(span.get("flags", 0) & 16)
        font.italic = bool(span.get("flags", 0) & 2)
        font.color.rgb = rgb_from_int(span.get("color", 0))


def add_image_block(slide, block):
    bbox = block.get("bbox")
    image = block.get("image")
    if not bbox or not image:
        return
    x0, y0, x1, y1 = bbox
    if x1 <= x0 or y1 <= y0:
        return
    slide.shapes.add_picture(
        io.BytesIO(image),
        pdf_to_inches(x0),
        pdf_to_inches(y0),
        pdf_to_inches(x1 - x0),
        pdf_to_inches(y1 - y0),
    )


def is_almost_full_page(rect, page_rect):
    x0, y0, x1, y1 = rect
    return (
        abs(x0 - page_rect.x0) < 2
        and abs(y0 - page_rect.y0) < 2
        and abs(x1 - page_rect.x1) < 2
        and abs(y1 - page_rect.y1) < 2
    )


def main():
    doc = fitz.open(PDF_PATH)
    prs = Presentation()
    prs.slide_width = pdf_to_inches(doc[0].rect.width)
    prs.slide_height = pdf_to_inches(doc[0].rect.height)
    blank = prs.slide_layouts[6]

    for page in doc:
        slide = prs.slides.add_slide(blank)

        # Draw vector fills and lines before text. Complex paths are approximated
        # by their bounding rectangles when they carry a fill.
        for drawing in page.get_drawings():
            rect = drawing.get("rect")
            fill = rgb_from_pdf(drawing.get("fill"))
            stroke = rgb_from_pdf(drawing.get("color"))
            width = drawing.get("width") or 0
            dtype = drawing.get("type")
            if rect and fill is not None:
                # Preserve white page background but skip duplicate full-page white fills.
                if fill == RGBColor(255, 255, 255) and is_almost_full_page(rect, page.rect):
                    continue
                add_rect(slide, rect, fill=fill, line=stroke, width=width)
            elif dtype == "s":
                for item in drawing.get("items", []):
                    if item[0] == "l":
                        add_line(slide, item[1], item[2], stroke, width)
                    elif item[0] == "re":
                        add_rect(slide, item[1], fill=None, line=stroke, width=width)

        text_dict = page.get_text("dict")
        for block in text_dict["blocks"]:
            if block.get("type") == 1:
                add_image_block(slide, block)

        for block in text_dict["blocks"]:
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                add_text_line(slide, line)

    # Remove the default first empty slide if python-pptx created one.
    if len(prs.slides) > doc.page_count:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    prs.save(OUT_PATH)
    print(OUT_PATH)


if __name__ == "__main__":
    main()
