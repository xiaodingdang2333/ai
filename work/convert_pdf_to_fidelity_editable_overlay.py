from pathlib import Path
import io

import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls


PDF_PATH = Path(r"D:\ai\work\qly_dict_service.pdf")
OUT_ASCII = Path(r"D:\ai\work\qly_dict_service_fidelity_editable_overlay.pptx")


def pdf_to_inches(v):
    return Inches(v / 72.0)


def rgb_from_int(value):
    return RGBColor((value >> 16) & 255, (value >> 8) & 255, value & 255)


def make_run_transparent(run):
    rpr = run._r.get_or_add_rPr()
    # python-pptx writes solidFill after font.color.rgb is set. Replace alpha
    # explicitly so the text remains editable/selectable but does not alter the
    # visible high-fidelity PDF background.
    solid = rpr.find(".//a:solidFill", namespaces={"a": "http://schemas.openxmlformats.org/drawingml/2006/main"})
    if solid is not None:
        srgb = solid.find(".//a:srgbClr", namespaces={"a": "http://schemas.openxmlformats.org/drawingml/2006/main"})
        if srgb is not None:
            srgb.append(parse_xml(f'<a:alpha {nsdecls("a")} val="0"/>'))


def add_invisible_text_line(slide, line):
    spans = [s for s in line.get("spans", []) if s.get("text")]
    if not spans:
        return
    text = "".join(s["text"] for s in spans)
    if not text.strip():
        return

    x0, y0, x1, y1 = line["bbox"]
    tx = slide.shapes.add_textbox(
        pdf_to_inches(x0),
        pdf_to_inches(y0 - 1),
        pdf_to_inches(max(x1 - x0 + 6, 8)),
        pdf_to_inches(max(y1 - y0 + 4, 8)),
    )
    tx.name = "可编辑文字层-" + text[:20]
    tf = tx.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1

    for span in spans:
        run = p.add_run()
        run.text = span["text"]
        font = run.font
        font.name = "Microsoft YaHei"
        font.size = Pt(max(1, span.get("size", 10)))
        font.bold = bool(span.get("flags", 0) & 16)
        font.italic = bool(span.get("flags", 0) & 2)
        font.color.rgb = rgb_from_int(span.get("color", 0))
        make_run_transparent(run)


def main():
    doc = fitz.open(PDF_PATH)
    prs = Presentation()
    prs.slide_width = pdf_to_inches(doc[0].rect.width)
    prs.slide_height = pdf_to_inches(doc[0].rect.height)
    blank = prs.slide_layouts[6]
    matrix = fitz.Matrix(2, 2)

    for page in doc:
        slide = prs.slides.add_slide(blank)
        pix = page.get_pixmap(matrix=matrix, alpha=False)
        slide.shapes.add_picture(
            io.BytesIO(pix.tobytes("png")),
            0,
            0,
            prs.slide_width,
            prs.slide_height,
        ).name = "高保真PDF底图"

        text_dict = page.get_text("dict")
        for block in text_dict["blocks"]:
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                add_invisible_text_line(slide, line)

    if len(prs.slides) > doc.page_count:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    prs.save(OUT_ASCII)
    print(OUT_ASCII)


if __name__ == "__main__":
    main()
