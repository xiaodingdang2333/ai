from __future__ import annotations

from pathlib import Path
import io
import json

import fitz
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


INPUT_PDF = Path(r"D:\ai\work\qly_dict_service.pdf")
OUTPUT_ASCII = Path(r"D:\ai\work\qly_dict_visible_editable.pptx")
OUTPUT_CN = Path(r"D:\ai\work\千里眼DICT集成服务细化管理方案宣贯材料_可见文字图形可编辑版.pptx")
REPORT_PATH = Path(r"D:\ai\work\qly_dict_visible_editable_report.json")
USE_TEXT_MASKS = False


def emu(v: float) -> int:
    return Inches(v / 72.0)


def rgb_from_int(value: int) -> RGBColor:
    return RGBColor((value >> 16) & 255, (value >> 8) & 255, value & 255)


def rgb_from_pdf(value) -> RGBColor:
    if value is None:
        return RGBColor(0, 0, 0)
    return RGBColor(*[max(0, min(255, int(round(c * 255)))) for c in value[:3]])


def render_page_white(page, scale=2):
    pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=True)
    rgba = Image.frombytes("RGBA", [pix.width, pix.height], pix.samples)
    white = Image.new("RGBA", rgba.size, (255, 255, 255, 255))
    white.alpha_composite(rgba)
    return white.convert("RGB")


def sample_background(img: Image.Image, bbox, scale: int) -> RGBColor:
    x0, y0, x1, y1 = bbox
    pts = [
        (x0 - 3, y0 - 3),
        (x0 - 3, (y0 + y1) / 2),
        (x1 + 3, y0 - 3),
        (x1 + 3, (y0 + y1) / 2),
        ((x0 + x1) / 2, y0 - 3),
        (x0 + 1, y0 + 1),
    ]
    colors = []
    for x, y in pts:
        px = int(max(0, min(img.width - 1, x * scale)))
        py = int(max(0, min(img.height - 1, y * scale)))
        colors.append(img.getpixel((px, py)))
    colors.sort(key=lambda c: sum(c))
    # Use a light-biased median; most text sits on white/light-blue fills.
    r, g, b = colors[len(colors) // 2]
    return RGBColor(r, g, b)


def add_mask(slide, bbox, fill: RGBColor, pad=1.2):
    x0, y0, x1, y1 = bbox
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        emu(x0 - pad),
        emu(y0 - pad),
        emu(max(x1 - x0 + pad * 2, 1)),
        emu(max(y1 - y0 + pad * 2, 1)),
    )
    shape.name = "TEXT_MASK"
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    shape.line.width = Pt(0.1)
    return shape


def add_text_line(slide, line):
    spans = [s for s in line.get("spans", []) if s.get("text")]
    if not spans:
        return False
    text = "".join(s["text"] for s in spans)
    if not text.strip():
        return False
    x0, y0, x1, y1 = line["bbox"]
    box = slide.shapes.add_textbox(
        emu(x0),
        emu(y0 - 1),
        emu(max(x1 - x0 + 8, 8)),
        emu(max(y1 - y0 + 5, 8)),
    )
    box.name = "EDITABLE_TEXT_" + text[:24]
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1
    for span in spans:
        run = p.add_run()
        run.text = span["text"]
        font = run.font
        font.name = "Microsoft YaHei"
        font.size = Pt(max(1, span.get("size", 10)))
        font.bold = bool(span.get("flags", 0) & 16)
        font.italic = bool(span.get("flags", 0) & 2)
        font.color.rgb = rgb_from_int(span.get("color", 0))
    return True


def add_simple_rect(slide, rect, fill=None, line=None, width=0.5, name="EDITABLE_SHAPE"):
    x0, y0, x1, y1 = rect
    if x1 <= x0 or y1 <= y0:
        return False
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(x0), emu(y0), emu(x1 - x0), emu(y1 - y0))
    shape.name = name
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is not None and width:
        shape.line.color.rgb = line
        shape.line.width = Pt(max(width, 0.25))
    else:
        shape.line.fill.background()
    return True


def add_line(slide, p1, p2, color, width):
    shape = slide.shapes.add_connector(1, emu(p1.x), emu(p1.y), emu(p2.x), emu(p2.y))
    shape.name = "EDITABLE_LINE"
    shape.line.color.rgb = color or RGBColor(0, 0, 0)
    shape.line.width = Pt(max(width or 0.5, 0.25))
    return True


def is_full_page(rect, page_rect) -> bool:
    x0, y0, x1, y1 = rect
    return abs(x0 - page_rect.x0) < 2 and abs(y0 - page_rect.y0) < 2 and abs(x1 - page_rect.x1) < 2 and abs(y1 - page_rect.y1) < 2


def add_visible_image(slide, block, page_rect):
    bbox = block.get("bbox")
    image = block.get("image")
    if not bbox or not image or is_full_page(bbox, page_rect):
        return False
    x0, y0, x1, y1 = bbox
    pic = slide.shapes.add_picture(io.BytesIO(image), emu(x0), emu(y0), emu(x1 - x0), emu(y1 - y0))
    pic.name = "EDITABLE_IMAGE"
    return True


def main():
    doc = fitz.open(INPUT_PDF)
    bg_doc = fitz.open(INPUT_PDF)
    for page in bg_doc:
        for block in page.get_text("dict")["blocks"]:
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                spans = [s for s in line.get("spans", []) if s.get("text")]
                if spans and "".join(s["text"] for s in spans).strip():
                    page.add_redact_annot(fitz.Rect(line["bbox"]))
        page.apply_redactions(images=0, graphics=0, text=0)
    prs = Presentation()
    prs.slide_width = emu(doc[0].rect.width)
    prs.slide_height = emu(doc[0].rect.height)
    blank = prs.slide_layouts[6]
    counts = {"pages": doc.page_count, "text": 0, "masks": 0, "images": 0, "rectangles": 0, "lines": 0}

    for page_index, page in enumerate(doc):
        slide = prs.slides.add_slide(blank)
        bg_img = render_page_white(bg_doc[page_index], scale=2)
        buf = io.BytesIO()
        bg_img.save(buf, format="PNG")
        buf.seek(0)
        slide.shapes.add_picture(buf, 0, 0, prs.slide_width, prs.slide_height).name = "FIDELITY_BACKGROUND"

        # Overlay editable simple rectangles only when they are represented by
        # simple PDF rectangle commands. Complex decorative paths and thousands
        # of segment lines stay in the background image to avoid visual damage.
        for drawing in page.get_drawings():
            fill = rgb_from_pdf(drawing.get("fill")) if drawing.get("fill") is not None else None
            stroke = rgb_from_pdf(drawing.get("color")) if drawing.get("color") is not None else None
            width = drawing.get("width") or 0.5
            for item in drawing.get("items", []):
                if item[0] == "re" and not is_full_page(item[1], page.rect):
                    if add_simple_rect(slide, item[1], fill=fill, line=stroke, width=width):
                        counts["rectangles"] += 1
                elif item[0] == "l":
                    continue

        text_dict = page.get_text("dict")
        # Do not overlay extracted PDF image blocks. Many are transparent-mask
        # artifacts with black backplates; the full-page background already
        # preserves the intended visuals.

        for block in text_dict["blocks"]:
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                spans = [s for s in line.get("spans", []) if s.get("text")]
                if not spans or not "".join(s["text"] for s in spans).strip():
                    continue
                if add_text_line(slide, line):
                    counts["text"] += 1

    if len(prs.slides) > doc.page_count:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]
    prs.save(OUTPUT_ASCII)
    OUTPUT_CN.write_bytes(OUTPUT_ASCII.read_bytes())
    REPORT_PATH.write_text(json.dumps(counts, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps(counts, ensure_ascii=False, indent=2))
    print(OUTPUT_CN)


if __name__ == "__main__":
    main()
