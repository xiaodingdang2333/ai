from pathlib import Path
import io

import fitz
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


PDF_PATH = Path(r"D:\ai\work\qly_dict_service.pdf")
OUT_ASCII = Path(r"D:\ai\work\qly_dict_service_hybrid_editable.pptx")
OUT_CN = Path(r"D:\ai\work\千里眼DICT集成服务细化管理方案宣贯材料_高保真可编辑混合版.pptx")


def pdf_to_inches(v):
    return Inches(v / 72.0)


def rgb_from_pdf(color):
    if color is None:
        return None
    return RGBColor(*[max(0, min(255, int(round(c * 255)))) for c in color[:3]])


def rgb_from_int(value):
    return RGBColor((value >> 16) & 255, (value >> 8) & 255, value & 255)


def add_rect(slide, rect, fill=None, line=None, width=0):
    x0, y0, x1, y1 = rect
    if x1 <= x0 or y1 <= y0:
        return None
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        pdf_to_inches(x0),
        pdf_to_inches(y0),
        pdf_to_inches(x1 - x0),
        pdf_to_inches(y1 - y0),
    )
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is not None and width and width > 0:
        shape.line.color.rgb = line
        shape.line.width = Pt(width)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, p1, p2, color, width):
    x0, y0 = p1
    x1, y1 = p2
    shape = slide.shapes.add_connector(
        1,
        pdf_to_inches(x0),
        pdf_to_inches(y0),
        pdf_to_inches(x1),
        pdf_to_inches(y1),
    )
    shape.line.color.rgb = color or RGBColor(0, 0, 0)
    shape.line.width = Pt(width or 0.5)
    return shape


def sample_bg_color(page_image, bbox, scale):
    x0, y0, x1, y1 = bbox
    cx = int(max(0, min(page_image.width - 1, ((x0 + x1) / 2) * scale)))
    cy = int(max(0, min(page_image.height - 1, ((y0 + y1) / 2) * scale)))
    r, g, b = page_image.getpixel((cx, cy))[:3]
    # If we sampled text itself, try the left edge of the text box.
    if r < 80 and g < 80 and b < 80:
        cx2 = int(max(0, min(page_image.width - 1, (x0 + 1) * scale)))
        cy2 = int(max(0, min(page_image.height - 1, (y0 + 1) * scale)))
        r, g, b = page_image.getpixel((cx2, cy2))[:3]
    return RGBColor(r, g, b)


def add_text_line(slide, line):
    spans = [s for s in line.get("spans", []) if s.get("text")]
    if not spans:
        return
    text = "".join(s["text"] for s in spans)
    if not text.strip():
        return

    x0, y0, x1, y1 = line["bbox"]
    height = max(y1 - y0 + 2, 8)
    width = max(x1 - x0 + 4, 8)
    tx = slide.shapes.add_textbox(
        pdf_to_inches(x0),
        pdf_to_inches(y0 - 1),
        pdf_to_inches(width),
        pdf_to_inches(height),
    )
    tf = tx.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1
    for span in spans:
        run = p.add_run()
        run.text = span["text"]
        font = run.font
        font.name = "Microsoft YaHei"
        font.size = Pt(max(1, span.get("size", 10)))
        font.bold = bool(span.get("flags", 0) & 16)
        font.italic = bool(span.get("flags", 0) & 2)
        font.color.rgb = rgb_from_int(span.get("color", 0))


def add_image_block(slide, block):
    bbox = block.get("bbox")
    image = block.get("image")
    if not bbox or not image:
        return
    x0, y0, x1, y1 = bbox
    if x1 <= x0 or y1 <= y0:
        return
    slide.shapes.add_picture(
        io.BytesIO(image),
        pdf_to_inches(x0),
        pdf_to_inches(y0),
        pdf_to_inches(x1 - x0),
        pdf_to_inches(y1 - y0),
    )


def is_full_page(rect, page_rect):
    x0, y0, x1, y1 = rect
    return (
        abs(x0 - page_rect.x0) < 2
        and abs(y0 - page_rect.y0) < 2
        and abs(x1 - page_rect.x1) < 2
        and abs(y1 - page_rect.y1) < 2
    )


def main():
    doc = fitz.open(PDF_PATH)
    prs = Presentation()
    prs.slide_width = pdf_to_inches(doc[0].rect.width)
    prs.slide_height = pdf_to_inches(doc[0].rect.height)
    blank = prs.slide_layouts[6]

    render_scale = 2
    matrix = fitz.Matrix(render_scale, render_scale)

    for page in doc:
        slide = prs.slides.add_slide(blank)
        pix = page.get_pixmap(matrix=matrix, alpha=False)
        bg_png = pix.tobytes("png")
        bg_img = Image.open(io.BytesIO(bg_png)).convert("RGB")
        slide.shapes.add_picture(io.BytesIO(bg_png), 0, 0, prs.slide_width, prs.slide_height)

        # Overlay simple editable shapes. The full-page bitmap remains the fidelity fallback.
        for drawing in page.get_drawings():
            rect = drawing.get("rect")
            fill = rgb_from_pdf(drawing.get("fill"))
            stroke = rgb_from_pdf(drawing.get("color"))
            width = drawing.get("width") or 0
            dtype = drawing.get("type")
            if rect and fill is not None and not is_full_page(rect, page.rect):
                add_rect(slide, rect, fill=fill, line=stroke, width=width)
            elif dtype == "s":
                for item in drawing.get("items", []):
                    if item[0] == "l":
                        add_line(slide, item[1], item[2], stroke, width)
                    elif item[0] == "re":
                        add_rect(slide, item[1], fill=None, line=stroke, width=width)

        text_dict = page.get_text("dict")

        # Add editable images over the rendered background.
        for block in text_dict["blocks"]:
            if block.get("type") == 1:
                add_image_block(slide, block)

        # Mask original rendered text line-by-line, then place editable text on top.
        for block in text_dict["blocks"]:
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                spans = [s for s in line.get("spans", []) if s.get("text")]
                if not spans or not "".join(s["text"] for s in spans).strip():
                    continue
                x0, y0, x1, y1 = line["bbox"]
                bg = sample_bg_color(bg_img, line["bbox"], render_scale)
                add_rect(slide, (x0 - 0.5, y0 - 0.5, x1 + 0.8, y1 + 0.8), fill=bg)
                add_text_line(slide, line)

    if len(prs.slides) > doc.page_count:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    prs.save(OUT_ASCII)
    OUT_CN.write_bytes(OUT_ASCII.read_bytes())
    print(OUT_CN)


if __name__ == "__main__":
    main()
