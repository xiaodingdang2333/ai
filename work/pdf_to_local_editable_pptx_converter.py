from __future__ import annotations

from pathlib import Path
import io
import json
import zipfile

import fitz
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls


INPUT_PDF = Path(r"D:\ai\work\qly_dict_service.pdf")
OUTPUT_ASCII = Path(r"D:\ai\work\qly_dict_local_high_fidelity_editable.pptx")
OUTPUT_CN = Path(r"D:\ai\work\千里眼DICT集成服务细化管理方案宣贯材料_本地高保真可编辑版.pptx")
REPORT_PATH = Path(r"D:\ai\work\千里眼DICT集成服务细化管理方案宣贯材料_本地转换验证报告.json")


def emu_from_pdf_points(value: float) -> int:
    return Inches(value / 72.0)


def rgb_from_int(value: int) -> RGBColor:
    return RGBColor((value >> 16) & 255, (value >> 8) & 255, value & 255)


def rgb_from_pdf(value) -> RGBColor:
    if value is None:
        return RGBColor(0, 0, 0)
    return RGBColor(*[max(0, min(255, int(round(c * 255)))) for c in value[:3]])


def append_alpha_zero_to_srgb(parent):
    srgb = parent.find(
        ".//a:srgbClr",
        namespaces={"a": "http://schemas.openxmlformats.org/drawingml/2006/main"},
    )
    if srgb is not None:
        srgb.append(parse_xml(f'<a:alpha {nsdecls("a")} val="0"/>'))


def make_run_transparent(run):
    rpr = run._r.get_or_add_rPr()
    append_alpha_zero_to_srgb(rpr)


def make_shape_transparent(shape):
    sp_pr = shape._element.spPr
    for solid in sp_pr.findall(
        ".//a:solidFill",
        namespaces={"a": "http://schemas.openxmlformats.org/drawingml/2006/main"},
    ):
        append_alpha_zero_to_srgb(solid)


def make_picture_transparent(pic):
    blip = pic._element.blipFill.blip
    # Keep the image editable/selectable while making it non-rendering.
    blip.append(parse_xml(f'<a:alphaModFix {nsdecls("a")} amt="0"/>'))


def add_invisible_text_line(slide, line):
    spans = [s for s in line.get("spans", []) if s.get("text")]
    if not spans:
        return False
    text = "".join(s["text"] for s in spans)
    if not text.strip():
        return False
    x0, y0, x1, y1 = line["bbox"]
    tx = slide.shapes.add_textbox(
        emu_from_pdf_points(x0),
        emu_from_pdf_points(y0 - 1),
        emu_from_pdf_points(max(x1 - x0 + 6, 8)),
        emu_from_pdf_points(max(y1 - y0 + 4, 8)),
    )
    tx.name = "EDIT_TEXT_" + text[:24]
    tf = tx.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1
    for span in spans:
        run = p.add_run()
        run.text = span["text"]
        font = run.font
        font.name = "Microsoft YaHei"
        font.size = Pt(max(1, span.get("size", 10)))
        font.bold = bool(span.get("flags", 0) & 16)
        font.italic = bool(span.get("flags", 0) & 2)
        font.color.rgb = rgb_from_int(span.get("color", 0))
        make_run_transparent(run)
    return True


def add_invisible_picture(slide, block):
    image = block.get("image")
    bbox = block.get("bbox")
    if not image or not bbox:
        return False
    x0, y0, x1, y1 = bbox
    if x1 <= x0 or y1 <= y0:
        return False
    pic = slide.shapes.add_picture(
        io.BytesIO(image),
        emu_from_pdf_points(x0),
        emu_from_pdf_points(y0),
        emu_from_pdf_points(x1 - x0),
        emu_from_pdf_points(y1 - y0),
    )
    pic.name = "EDIT_IMAGE"
    make_picture_transparent(pic)
    return True


def add_invisible_rect(slide, rect, name="EDIT_SHAPE"):
    x0, y0, x1, y1 = rect
    if x1 <= x0 or y1 <= y0:
        return False
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        emu_from_pdf_points(x0),
        emu_from_pdf_points(y0),
        emu_from_pdf_points(x1 - x0),
        emu_from_pdf_points(y1 - y0),
    )
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = RGBColor(255, 255, 255)
    shape.line.width = Pt(0.25)
    make_shape_transparent(shape)
    return True


def add_invisible_line(slide, p1, p2):
    x0, y0 = p1
    x1, y1 = p2
    shape = slide.shapes.add_connector(
        1,
        emu_from_pdf_points(x0),
        emu_from_pdf_points(y0),
        emu_from_pdf_points(x1),
        emu_from_pdf_points(y1),
    )
    shape.name = "EDIT_LINE"
    shape.line.color.rgb = RGBColor(255, 255, 255)
    shape.line.width = Pt(0.25)
    make_shape_transparent(shape)
    return True


def is_full_page(rect, page_rect) -> bool:
    x0, y0, x1, y1 = rect
    return (
        abs(x0 - page_rect.x0) < 2
        and abs(y0 - page_rect.y0) < 2
        and abs(x1 - page_rect.x1) < 2
        and abs(y1 - page_rect.y1) < 2
    )


def build_pptx() -> dict:
    doc = fitz.open(INPUT_PDF)
    prs = Presentation()
    prs.slide_width = emu_from_pdf_points(doc[0].rect.width)
    prs.slide_height = emu_from_pdf_points(doc[0].rect.height)
    blank = prs.slide_layouts[6]

    counts = {
        "pages": doc.page_count,
        "backgrounds": 0,
        "editable_text_lines": 0,
        "editable_images": 0,
        "editable_shapes": 0,
        "editable_lines": 0,
    }

    for page in doc:
        slide = prs.slides.add_slide(blank)
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=True)
        rgba = Image.frombytes("RGBA", [pix.width, pix.height], pix.samples)
        white = Image.new("RGBA", rgba.size, (255, 255, 255, 255))
        white.alpha_composite(rgba)
        png_buffer = io.BytesIO()
        white.convert("RGB").save(png_buffer, format="PNG")
        png_buffer.seek(0)
        bg = slide.shapes.add_picture(
            png_buffer,
            0,
            0,
            prs.slide_width,
            prs.slide_height,
        )
        bg.name = "VISIBLE_PDF_BACKGROUND"
        counts["backgrounds"] += 1

        text_dict = page.get_text("dict")
        for block in text_dict["blocks"]:
            if block.get("type") == 1 and add_invisible_picture(slide, block):
                counts["editable_images"] += 1

        for drawing in page.get_drawings():
            rect = drawing.get("rect")
            if rect and not is_full_page(rect, page.rect):
                # Treat filled vector regions as editable bounding rectangles.
                if drawing.get("fill") is not None and add_invisible_rect(slide, rect):
                    counts["editable_shapes"] += 1
            if drawing.get("type") == "s":
                for item in drawing.get("items", []):
                    if item[0] == "l" and add_invisible_line(slide, item[1], item[2]):
                        counts["editable_lines"] += 1
                    elif item[0] == "re" and add_invisible_rect(slide, item[1]):
                        counts["editable_shapes"] += 1

        for block in text_dict["blocks"]:
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                if add_invisible_text_line(slide, line):
                    counts["editable_text_lines"] += 1

    if len(prs.slides) > doc.page_count:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    prs.save(OUTPUT_ASCII)
    OUTPUT_CN.write_bytes(OUTPUT_ASCII.read_bytes())
    return counts


def inspect_pptx_package() -> dict:
    with zipfile.ZipFile(OUTPUT_ASCII, "r") as zf:
        slide_xmls = [n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        alpha_zero = 0
        bg_count = 0
        editable_named = 0
        for name in slide_xmls:
            xml = zf.read(name).decode("utf-8", errors="ignore")
            alpha_zero += xml.count('val="0"') + xml.count('amt="0"')
            bg_count += xml.count("VISIBLE_PDF_BACKGROUND")
            editable_named += xml.count("EDIT_TEXT_") + xml.count("EDIT_IMAGE") + xml.count("EDIT_SHAPE") + xml.count("EDIT_LINE")
    return {
        "slide_xml_count": len(slide_xmls),
        "visible_background_objects": bg_count,
        "transparent_alpha_markers": alpha_zero,
        "editable_named_objects": editable_named,
    }


if __name__ == "__main__":
    result = {"build_counts": build_pptx(), "package_inspection": inspect_pptx_package()}
    REPORT_PATH.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps(result, ensure_ascii=False, indent=2))
    print(OUTPUT_CN)
